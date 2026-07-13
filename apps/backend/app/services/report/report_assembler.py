"""
ReportAssembler — Step 10 AI Report Generator

Multi-format export engine. Receives the LLM narrative (markdown) + collected
structured data and assembles the final file.

Supported formats:
  - markdown  : Raw .md file (always generated as source of truth)
  - pdf       : PDF via ReportLab (fallback: save .txt if unavailable)
  - excel     : XLSX via openpyxl (stats table + narrative sheet)
  - pptx      : PowerPoint via python-pptx (title + bullet slides)
  - png       : PNG summary card via Pillow

All output files are saved under: storage/reports/{report_id}/

Design principle: each export method is independently callable so future
AI Agents can request specific formats without re-running the full pipeline.
"""
import os
import logging
import textwrap
from typing import Dict, Any

logger = logging.getLogger("app.services.report.report_assembler")

REPORT_DIR = os.path.join("storage", "reports")
os.makedirs(REPORT_DIR, exist_ok=True)


class ReportAssembler:

    @classmethod
    def export(
        cls,
        report_id: str,
        narrative: str,
        report_type: str,
        export_format: str,
        collected: Dict[str, Any],
    ) -> str:
        """
        Routes to the correct exporter based on export_format.
        Always writes the .md source file first.
        Returns the absolute path of the primary exported file.
        """
        report_subdir = os.path.join(REPORT_DIR, report_id)
        os.makedirs(report_subdir, exist_ok=True)

        # Always write markdown source
        md_path = cls._write_markdown(report_subdir, report_id, narrative)

        fmt = export_format.lower()
        if fmt == "pdf":
            return cls._export_pdf(report_subdir, report_id, narrative, collected)
        elif fmt == "excel":
            return cls._export_excel(report_subdir, report_id, narrative, collected)
        elif fmt == "pptx":
            return cls._export_pptx(report_subdir, report_id, narrative, collected)
        elif fmt == "png":
            return cls._export_png(report_subdir, report_id, narrative, collected)
        else:
            # Default: return markdown
            return md_path

    # ──────────────────────────────────────────────────────────────────
    # Markdown (source of truth)
    # ──────────────────────────────────────────────────────────────────

    @classmethod
    def _write_markdown(cls, subdir: str, report_id: str, narrative: str) -> str:
        path = os.path.join(subdir, f"{report_id}.md")
        with open(path, "w", encoding="utf-8") as f:
            f.write(narrative)
        logger.info(f"[Assembler] Markdown written: {path}")
        return path

    # ──────────────────────────────────────────────────────────────────
    # PDF via ReportLab
    # ──────────────────────────────────────────────────────────────────

    @classmethod
    def _export_pdf(
        cls, subdir: str, report_id: str, narrative: str, collected: Dict[str, Any]
    ) -> str:
        path = os.path.join(subdir, f"{report_id}.pdf")
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import cm
            from reportlab.lib import colors
            from reportlab.platypus import (
                SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
            )

            doc = SimpleDocTemplate(
                path,
                pagesize=A4,
                leftMargin=2 * cm,
                rightMargin=2 * cm,
                topMargin=2 * cm,
                bottomMargin=2 * cm,
            )

            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                "ReportTitle",
                parent=styles["Title"],
                fontSize=20,
                textColor=colors.HexColor("#1a1a2e"),
                spaceAfter=12,
            )
            heading_style = ParagraphStyle(
                "SectionHeading",
                parent=styles["Heading2"],
                fontSize=13,
                textColor=colors.HexColor("#4a4de7"),
                spaceBefore=16,
                spaceAfter=6,
            )
            body_style = ParagraphStyle(
                "BodyText",
                parent=styles["BodyText"],
                fontSize=10,
                leading=16,
                textColor=colors.HexColor("#222222"),
            )

            story = []

            # Title
            doc_id = collected.get("doc_id", "?")
            story.append(Paragraph(f"Nexora AI Report — Document {doc_id}", title_style))
            story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#4a4de7")))
            story.append(Spacer(1, 0.4 * cm))

            # Narrative — parse markdown headings into PDF paragraphs
            for line in narrative.split("\n"):
                line_stripped = line.strip()
                if not line_stripped:
                    story.append(Spacer(1, 0.2 * cm))
                elif line_stripped.startswith("## "):
                    story.append(Paragraph(line_stripped[3:], heading_style))
                elif line_stripped.startswith("# "):
                    story.append(Paragraph(line_stripped[2:], title_style))
                elif line_stripped.startswith("### "):
                    story.append(Paragraph(line_stripped[4:], heading_style))
                elif line_stripped.startswith("- ") or line_stripped.startswith("* "):
                    story.append(Paragraph(f"• {line_stripped[2:]}", body_style))
                elif line_stripped.startswith("|"):
                    # Skip markdown table rows in PDF (tables added separately)
                    pass
                else:
                    story.append(Paragraph(line_stripped, body_style))

            # ML Metrics Table (if available)
            ml = collected.get("ml_session", {})
            if ml.get("model_trained") and ml.get("metrics"):
                story.append(Spacer(1, 0.5 * cm))
                story.append(Paragraph("Model Performance Metrics", heading_style))
                table_data = [["Metric", "Value"]] + [
                    [k, str(v)] for k, v in ml["metrics"].items()
                ]
                t = Table(table_data, colWidths=[8 * cm, 6 * cm])
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4a4de7")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#f5f5ff"), colors.white]),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
                    ("ALIGN", (1, 0), (1, -1), "CENTER"),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ]))
                story.append(t)

            # SHAP Table (if available)
            shap = collected.get("shap", {})
            if shap.get("available") and shap.get("summary"):
                story.append(Spacer(1, 0.5 * cm))
                story.append(Paragraph("SHAP Feature Attributions", heading_style))
                shap_data = [["Feature", "Mean |SHAP|"]] + [
                    [item["feature"], f"{item['mean_abs_shap']:.4f}"]
                    for item in shap["summary"][:10]
                ]
                st = Table(shap_data, colWidths=[10 * cm, 4 * cm])
                st.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1a1a2e")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#f0f0ff"), colors.white]),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
                    ("ALIGN", (1, 0), (1, -1), "CENTER"),
                    ("TOPPADDING", (0, 0), (-1, -1), 5),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ]))
                story.append(st)

            doc.build(story)
            logger.info(f"[Assembler] PDF written: {path}")
            return path

        except ImportError:
            logger.warning("[Assembler] reportlab not installed — saving .txt fallback")
            txt_path = path.replace(".pdf", "_fallback.txt")
            with open(txt_path, "w", encoding="utf-8") as f:
                f.write(narrative)
            return txt_path
        except Exception as e:
            logger.error(f"[Assembler] PDF generation failed: {e}")
            # Return markdown path as safe fallback
            return os.path.join(subdir, f"{report_id}.md")

    # ──────────────────────────────────────────────────────────────────
    # Excel via openpyxl
    # ──────────────────────────────────────────────────────────────────

    @classmethod
    def _export_excel(
        cls, subdir: str, report_id: str, narrative: str, collected: Dict[str, Any]
    ) -> str:
        path = os.path.join(subdir, f"{report_id}.xlsx")
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter

            wb = openpyxl.Workbook()

            # ── Sheet 1: Narrative ──────────────────────────────────
            ws_narrative = wb.active
            ws_narrative.title = "AI Narrative"
            ws_narrative.column_dimensions["A"].width = 120
            header_fill = PatternFill(start_color="1a1a2e", end_color="1a1a2e", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF", size=12)
            ws_narrative["A1"] = "AI-Generated Report Narrative"
            ws_narrative["A1"].font = header_font
            ws_narrative["A1"].fill = header_fill
            ws_narrative["A1"].alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
            ws_narrative.row_dimensions[1].height = 25

            for i, line in enumerate(narrative.split("\n"), start=2):
                cell = ws_narrative.cell(row=i, column=1, value=line)
                cell.alignment = Alignment(wrap_text=True)
                if line.startswith("# ") or line.startswith("## "):
                    cell.font = Font(bold=True, color="4a4de7", size=11)
                elif line.startswith("- ") or line.startswith("* "):
                    cell.value = f"• {line[2:]}"

            # ── Sheet 2: ML Metrics ─────────────────────────────────
            ml = collected.get("ml_session", {})
            if ml.get("model_trained"):
                ws_ml = wb.create_sheet("ML Metrics")
                headers = ["Metric", "Value"]
                for col_idx, header in enumerate(headers, 1):
                    cell = ws_ml.cell(row=1, column=col_idx, value=header)
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill(start_color="4a4de7", end_color="4a4de7", fill_type="solid")
                    cell.alignment = Alignment(horizontal="center")

                for row_idx, (metric, value) in enumerate(ml.get("metrics", {}).items(), start=2):
                    ws_ml.cell(row=row_idx, column=1, value=metric)
                    ws_ml.cell(row=row_idx, column=2, value=value)
                    fill_color = "f5f5ff" if row_idx % 2 == 0 else "ffffff"
                    for c in range(1, 3):
                        ws_ml.cell(row=row_idx, column=c).fill = PatternFill(
                            start_color=fill_color, end_color=fill_color, fill_type="solid"
                        )

                # Feature Importance
                fi = ml.get("feature_importance", [])
                if fi:
                    start_row = len(ml.get("metrics", {})) + 4
                    ws_ml.cell(row=start_row, column=1, value="Feature Importance").font = Font(bold=True)
                    for i, item in enumerate(fi[:20], start=start_row + 1):
                        ws_ml.cell(row=i, column=1, value=item["feature"])
                        ws_ml.cell(row=i, column=2, value=item["importance"])

                ws_ml.column_dimensions["A"].width = 30
                ws_ml.column_dimensions["B"].width = 20

            # ── Sheet 3: SHAP Attributions ──────────────────────────
            shap = collected.get("shap", {})
            if shap.get("available") and shap.get("summary"):
                ws_shap = wb.create_sheet("SHAP Explanation")
                for col_idx, header in enumerate(["Feature", "Mean |SHAP|"], 1):
                    cell = ws_shap.cell(row=1, column=col_idx, value=header)
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill(start_color="1a1a2e", end_color="1a1a2e", fill_type="solid")
                    cell.alignment = Alignment(horizontal="center")
                for i, item in enumerate(shap["summary"], start=2):
                    ws_shap.cell(row=i, column=1, value=item["feature"])
                    ws_shap.cell(row=i, column=2, value=item["mean_abs_shap"])
                ws_shap.column_dimensions["A"].width = 30
                ws_shap.column_dimensions["B"].width = 20

            # ── Sheet 4: Analytics Stats ────────────────────────────
            analytics = collected.get("analytics", {})
            col_stats = analytics.get("column_stats", {})
            if col_stats:
                ws_stats = wb.create_sheet("Analytics Stats")
                stat_headers = ["Column", "Type", "Mean", "Std", "Min", "Max", "Missing"]
                for col_idx, header in enumerate(stat_headers, 1):
                    cell = ws_stats.cell(row=1, column=col_idx, value=header)
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill(start_color="4a4de7", end_color="4a4de7", fill_type="solid")
                    cell.alignment = Alignment(horizontal="center")
                    ws_stats.column_dimensions[get_column_letter(col_idx)].width = 16

                for row_idx, (col, stats) in enumerate(col_stats.items(), start=2):
                    ws_stats.cell(row=row_idx, column=1, value=col)
                    ws_stats.cell(row=row_idx, column=2, value=stats.get("dtype", "?"))
                    ws_stats.cell(row=row_idx, column=3, value=stats.get("mean", ""))
                    ws_stats.cell(row=row_idx, column=4, value=stats.get("std", ""))
                    ws_stats.cell(row=row_idx, column=5, value=stats.get("min", ""))
                    ws_stats.cell(row=row_idx, column=6, value=stats.get("max", ""))
                    ws_stats.cell(row=row_idx, column=7, value=stats.get("missing", 0))

            wb.save(path)
            logger.info(f"[Assembler] Excel written: {path}")
            return path

        except Exception as e:
            logger.error(f"[Assembler] Excel export failed: {e}")
            return os.path.join(subdir, f"{report_id}.md")

    # ──────────────────────────────────────────────────────────────────
    # PowerPoint via python-pptx
    # ──────────────────────────────────────────────────────────────────

    @classmethod
    def _export_pptx(
        cls, subdir: str, report_id: str, narrative: str, collected: Dict[str, Any]
    ) -> str:
        path = os.path.join(subdir, f"{report_id}.pptx")
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            DARK = RGBColor(0x1a, 0x1a, 0x2e)
            ACCENT = RGBColor(0x4a, 0x4d, 0xe7)
            WHITE = RGBColor(0xFF, 0xFF, 0xFF)
            LIGHT = RGBColor(0xF5, 0xF5, 0xFF)

            blank_layout = prs.slide_layouts[6]  # blank

            def add_slide_with_title(title_text: str, body_lines: list, bg_color=DARK):
                slide = prs.slides.add_slide(blank_layout)
                # Background
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = bg_color

                # Title box
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.1))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = WHITE

                # Accent line
                from pptx.util import Pt as PtUnit
                connector = slide.shapes.add_shape(
                    1,  # MSO_SHAPE_TYPE.RECTANGLE
                    Inches(0.5), Inches(1.5), Inches(12), Inches(0.05)
                )
                connector.fill.solid()
                connector.fill.fore_color.rgb = ACCENT
                connector.line.fill.background()

                # Body text
                body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5))
                btf = body_box.text_frame
                btf.word_wrap = True
                for i, line in enumerate(body_lines[:18]):
                    if i == 0:
                        bp = btf.paragraphs[0]
                    else:
                        bp = btf.add_paragraph()
                    bp.text = line if line else " "
                    bp.font.size = Pt(13)
                    bp.font.color.rgb = LIGHT if bg_color == DARK else DARK
                    if line.startswith("## ") or line.startswith("# "):
                        bp.text = line.lstrip("#").strip()
                        bp.font.bold = True
                        bp.font.size = Pt(15)
                        bp.font.color.rgb = ACCENT

                return slide

            # ── Title Slide ─────────────────────────────────────────
            doc_id = collected.get("doc_id", "?")
            slide = prs.slides.add_slide(blank_layout)
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = DARK

            title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Nexora AI Report"
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = WHITE

            sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10), Inches(1))
            stf = sub_box.text_frame
            sp = stf.paragraphs[0]
            sp.text = f"Document ID: {doc_id}  |  Generated by Nexora AI"
            sp.font.size = Pt(16)
            sp.font.color.rgb = ACCENT

            # ── Narrative Slides (chunked) ──────────────────────────
            lines = narrative.split("\n")
            chunks = [lines[i:i + 16] for i in range(0, len(lines), 16)]
            for idx, chunk in enumerate(chunks[:12]):  # max 12 content slides
                title = next((line.lstrip("#").strip() for line in chunk if line.strip().startswith("#")), f"Report Section {idx + 1}")
                body = [line for line in chunk if not line.strip().startswith("#")]
                add_slide_with_title(title, body)

            # ── ML Metrics Slide ────────────────────────────────────
            ml = collected.get("ml_session", {})
            if ml.get("model_trained") and ml.get("metrics"):
                metrics_lines = [f"Algorithm: {ml.get('algorithm', '?')}",
                                 f"Task: {ml.get('task_type', '?')}", ""]
                metrics_lines += [f"• {k}: {v}" for k, v in ml["metrics"].items()]
                add_slide_with_title("Model Performance", metrics_lines, bg_color=RGBColor(0x0d, 0x0d, 0x1a))

            # ── SHAP Slide ──────────────────────────────────────────
            shap = collected.get("shap", {})
            if shap.get("available") and shap.get("summary"):
                shap_lines = ["Top Feature Attributions (Mean |SHAP|)", ""]
                shap_lines += [f"• {item['feature']}: {item['mean_abs_shap']:.4f}" for item in shap["summary"][:10]]
                add_slide_with_title("SHAP Explainability", shap_lines, bg_color=RGBColor(0x0d, 0x0d, 0x1a))

            prs.save(path)
            logger.info(f"[Assembler] PPTX written: {path}")
            return path

        except ImportError:
            logger.warning("[Assembler] python-pptx not installed — returning markdown")
            return os.path.join(subdir, f"{report_id}.md")
        except Exception as e:
            logger.error(f"[Assembler] PPTX export failed: {e}")
            return os.path.join(subdir, f"{report_id}.md")

    # ──────────────────────────────────────────────────────────────────
    # PNG Summary Card via Pillow
    # ──────────────────────────────────────────────────────────────────

    @classmethod
    def _export_png(
        cls, subdir: str, report_id: str, narrative: str, collected: Dict[str, Any]
    ) -> str:
        path = os.path.join(subdir, f"{report_id}.png")
        try:
            from PIL import Image, ImageDraw, ImageFont

            W, H = 1600, 900
            DARK_BG = (26, 26, 46)
            ACCENT = (74, 77, 231)
            WHITE = (255, 255, 255)
            LIGHT_TEXT = (220, 220, 240)
            MUTED = (140, 140, 160)

            img = Image.new("RGB", (W, H), color=DARK_BG)
            draw = ImageDraw.Draw(img)

            # Header bar
            draw.rectangle([(0, 0), (W, 90)], fill=ACCENT)
            draw.text((40, 22), "Nexora AI Report", fill=WHITE,
                      font=cls._pil_font(36, bold=True))
            doc_id = collected.get("doc_id", "?")
            draw.text((40, 60), f"Document ID: {doc_id}", fill=WHITE,
                      font=cls._pil_font(18))

            # Narrative preview (first 30 lines)
            y = 110
            lines = narrative.split("\n")[:35]
            for line in lines:
                if not line.strip():
                    y += 12
                    continue
                color = ACCENT if line.startswith("#") else LIGHT_TEXT
                size = 22 if line.startswith("#") else 16
                wrapped = textwrap.wrap(line.lstrip("#").strip(), width=100)
                for wline in wrapped[:2]:
                    if y > H - 120:
                        break
                    draw.text((40, y), wline, fill=color, font=cls._pil_font(size))
                    y += size + 8
                if y > H - 120:
                    break

            # ML metrics panel (right side)
            ml = collected.get("ml_session", {})
            if ml.get("model_trained") and ml.get("metrics"):
                draw.rectangle([(W - 380, 95), (W - 20, H - 20)], fill=(30, 30, 55))
                draw.text((W - 370, 105), "Model Metrics", fill=ACCENT,
                          font=cls._pil_font(20, bold=True))
                draw.text((W - 370, 130), ml.get("algorithm", "?"), fill=MUTED,
                          font=cls._pil_font(14))
                my = 160
                for k, v in list(ml["metrics"].items())[:6]:
                    draw.text((W - 370, my), f"{k}:", fill=MUTED, font=cls._pil_font(13))
                    draw.text((W - 200, my), str(v), fill=WHITE, font=cls._pil_font(13, bold=True))
                    my += 28

                # Feature importance mini bars
                fi = ml.get("feature_importance", [])[:5]
                if fi:
                    my += 10
                    draw.text((W - 370, my), "Top Features", fill=ACCENT, font=cls._pil_font(16, bold=True))
                    my += 24
                    max_fi = fi[0]["importance"] if fi else 1
                    for item in fi:
                        bar_w = int((item["importance"] / max_fi) * 300)
                        draw.rectangle([(W - 370, my + 4), (W - 370 + bar_w, my + 16)], fill=ACCENT)
                        draw.text((W - 370, my + 20), item["feature"][:22], fill=LIGHT_TEXT, font=cls._pil_font(12))
                        my += 46

            # SHAP indicator badge
            shap = collected.get("shap", {})
            badge_color = (50, 200, 100) if shap.get("available") else (180, 60, 60)
            badge_text = "✓ SHAP Available" if shap.get("available") else "○ SHAP Not Available"
            draw.rectangle([(40, H - 60), (300, H - 20)], fill=badge_color)
            draw.text((55, H - 50), badge_text, fill=WHITE, font=cls._pil_font(14, bold=True))

            # Footer
            draw.rectangle([(0, H - 15), (W, H)], fill=ACCENT)

            img.save(path, "PNG")
            logger.info(f"[Assembler] PNG written: {path}")
            return path

        except ImportError:
            logger.warning("[Assembler] Pillow not installed — returning markdown")
            return os.path.join(subdir, f"{report_id}.md")
        except Exception as e:
            logger.error(f"[Assembler] PNG export failed: {e}")
            return os.path.join(subdir, f"{report_id}.md")

    @staticmethod
    def _pil_font(size: int, bold: bool = False):
        """Loads a PIL font, falling back to default if fonts unavailable."""
        try:
            from PIL import ImageFont
            # Try loading a system font
            font_candidates = [
                "arial.ttf", "Arial.ttf", "DejaVuSans.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "C:/Windows/Fonts/arial.ttf",
            ]
            for candidate in font_candidates:
                try:
                    return ImageFont.truetype(candidate, size)
                except Exception:
                    continue
            return ImageFont.load_default()
        except Exception:
            from PIL import ImageFont
            return ImageFont.load_default()
